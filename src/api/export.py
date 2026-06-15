import base64
import json
import os
import re
import shutil
import subprocess
import tempfile
import unicodedata
import zipfile
from io import BytesIO

from fastapi import APIRouter, BackgroundTasks, Depends, HTTPException, status
from fastapi.responses import FileResponse, HTMLResponse, PlainTextResponse
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pydantic import BaseModel
from sqlalchemy.orm import Session

from src.auth import get_current_user
from src.database.models import CLO, Chapter, ChapterMaterial, Course, Question, User
from src.database.session import get_db
from src.services.slide_renderer import render_slide_to_svg
from src.utils.active_learning_parser import parse_active_learning_into_notes
from src.utils.markdown_to_slidej import parse_markdown_to_slides

router = APIRouter(prefix="/api/courses", tags=["export"])


@router.get("/{course_id}/export-materials", response_class=PlainTextResponse)
def export_course_materials(
    course_id: int, current_user: User = Depends(get_current_user), db: Session = Depends(get_db)
):
    # 1. Xác thực môn học
    course = db.query(Course).filter(Course.id == course_id, Course.user_id == current_user.id).first()
    if not course:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND, detail="Môn học không tồn tại hoặc bạn không có quyền truy cập."
        )

    # 2. Lấy danh sách chương học sắp xếp theo sort_order
    chapters = (
        db.query(Chapter)
        .filter(Chapter.course_id == course_id, Chapter.is_active)
        .order_by(Chapter.sort_order)
        .all()
    )

    # 3. Tạo nội dung file Markdown tổng hợp
    content = f"# GIÁO ÁN HỌC LIỆU MÔN HỌC: {course.course_name.upper()}\n"
    content += f"Mã môn học: {course.course_code}\n"
    content += f"Giảng viên biên soạn: {current_user.full_name or current_user.email}\n"
    content += "Sinh tự động bởi AI Lecture Assistant (G02-Team023)\n\n"
    content += "========================================================\n\n"

    if not chapters:
        content += "* Chưa có nội dung chương học nào được thiết kế cho môn học này.\n"
    else:
        for idx, ch in enumerate(chapters):
            content += f"## CHƯƠNG {idx + 1}: {ch.title.upper()}\n"
            content += f"Mô tả chương: {ch.description or 'N/A'}\n\n"

            # Lấy học liệu của chương
            material = db.query(ChapterMaterial).filter(ChapterMaterial.chapter_id == ch.id).first()
            if material:
                content += "### 1. Slide Bài giảng (Markdown)\n"
                if material.slide_content:
                    content += f"{material.slide_content}\n\n"
                else:
                    content += "* Chưa biên soạn slide cho chương này.\n\n"

                content += "### 2. Kịch bản Hoạt động (Active Learning)\n"
                if material.active_learning_script:
                    content += f"{material.active_learning_script}\n\n"
                else:
                    content += "* Chưa biên soạn kịch bản active learning cho chương này.\n\n"
            else:
                content += "* Chương học chưa được thiết kế học liệu.\n\n"

            content += "--------------------------------------------------------\n\n"

    # Thiết lập headers để trình duyệt nhận diện tải file đính kèm
    headers = {"Content-Disposition": f"attachment; filename=Giao_an_{course.course_code}.md"}
    return PlainTextResponse(content, headers=headers)


@router.get("/{course_id}/export-questions", response_class=PlainTextResponse)
def export_course_questions(
    course_id: int, current_user: User = Depends(get_current_user), db: Session = Depends(get_db)
):
    # 1. Xác thực môn học
    course = db.query(Course).filter(Course.id == course_id, Course.user_id == current_user.id).first()
    if not course:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND, detail="Môn học không tồn tại hoặc bạn không có quyền truy cập."
        )

    # 2. Lấy danh sách câu hỏi
    questions = db.query(Question).filter(Question.course_id == course_id, Question.is_active).all()
    clos = db.query(CLO).filter(CLO.course_id == course_id).all()

    # 3. Tạo đề thi
    content = f"# ĐỀ THI TRẮC NGHIỆM MÔN HỌC: {course.course_name.upper()}\n"
    content += f"Mã môn học: {course.course_code}\n"
    content += f"Số lượng câu hỏi: {len(questions)} câu\n"
    content += "Thời gian làm bài: 45 phút (Đề thi tham khảo)\n"
    content += "========================================================\n\n"

    if not questions:
        content += "* Chưa soạn câu hỏi thi trắc nghiệm nào trong ngân hàng đề thi.\n"
    else:
        # Phần 1: Đề thi
        content += "## PHẦN I: ĐỀ THI\n\n"
        for idx, q in enumerate(questions):
            content += f"Câu {idx + 1}: {q.question_text}\n"

            opts = []
            try:
                opts = json.loads(q.options_json) if q.options_json else []
            except Exception:
                opts = []

            labels = ["A", "B", "C", "D"]
            for o_idx, opt in enumerate(opts):
                if o_idx < len(labels):
                    content += f"  {labels[o_idx]}. {opt}\n"
            content += "\n"

        # Phần 2: Đáp án đối chiếu
        content += "========================================================\n\n"
        content += "## PHẦN II: ĐÁP ÁN VÀ MA TRẬN PHÂN LOẠI CHẤT LƯỢNG (CLO - BLOOM)\n\n"

        for idx, q in enumerate(questions):
            linked_clo = next((c for c in clos if c.id == q.clo_id), None)
            clo_code = linked_clo.clo_code if linked_clo else "N/A"

            content += f"Câu {idx + 1}:\n"
            content += f"  - Đáp án đúng: {q.correct_answer}\n"
            content += f"  - Chuẩn đầu ra: {clo_code}\n"
            content += f"  - Cấp độ Bloom: Mức {q.bloom_level}\n\n"

    headers = {"Content-Disposition": f"attachment; filename=De_thi_{course.course_code}.md"}
    return PlainTextResponse(content, headers=headers)


# SlideJ logic removed


# PPT-Master script path: ưu tiên biến môi trường, fallback sang đường dẫn mặc định
_PPT_MASTER_SCRIPT_PATH = os.environ.get(
    "PPT_MASTER_SCRIPT_PATH",
    os.path.join(
        os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__)))),
        "..",
        "ContentToSlide",
        "ppt-master",
        "skills",
        "ppt-master",
        "scripts",
        "svg_to_pptx.py",
    ),
)


@router.get("/chapters/{chapter_id}/export-pptx")
def export_chapter_pptx(
    chapter_id: int,
    theme: str = "warm_academic",
    engine: str = "ppt_master",
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db),
):
    # 1. Xác thực chương học và quyền sở hữu môn học
    chapter = (
        db.query(Chapter)
        .join(Course)
        .filter(Chapter.id == chapter_id, Course.user_id == current_user.id, Chapter.is_active)
        .first()
    )
    if not chapter:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND, detail="Chương học không tồn tại hoặc bạn không có quyền truy cập."
        )

    # 2. Lấy học liệu chương học
    material = db.query(ChapterMaterial).filter(ChapterMaterial.chapter_id == chapter_id).first()
    if not material or not material.slide_content:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST, detail="Chương học này chưa có nội dung slide bài giảng thiết kế."
        )

    # Always use PPT-Master engine, SlideJ has been removed
    if True:
        import shutil
        import sys

        # 1. Parse slide markdown content
        parsed_slides = parse_markdown_to_slides(material.slide_content)

        # 2. Setup project directories
        project_dir = os.path.join("temp", f"ppt_master_{chapter_id}")
        svg_output_dir = os.path.join(project_dir, "svg_output")
        notes_dir = os.path.join(project_dir, "notes")
        exports_dir = os.path.join(project_dir, "exports")

        # Cleanup old run if exists
        if os.path.exists(project_dir):
            try:
                shutil.rmtree(project_dir)
            except Exception:
                pass

        os.makedirs(svg_output_dir, exist_ok=True)
        os.makedirs(notes_dir, exist_ok=True)
        os.makedirs(exports_dir, exist_ok=True)

        # Parse active learning notes map
        notes_by_slide = parse_active_learning_into_notes(material.active_learning_script)

        # 3. Create SVG and Notes files
        for idx, s in enumerate(parsed_slides):
            slide_name = f"slide_{idx + 1:05d}"
            svg_file_path = os.path.join(svg_output_dir, f"{slide_name}.svg")

            # Use inline SVG if available
            if s.get("svg_content"):
                with open(svg_file_path, "w", encoding="utf-8") as f:
                    f.write(s["svg_content"])
            else:
                # Dynamic generation of SVG slide using slide_renderer service
                svg_content = render_slide_to_svg(s, theme, idx)
                with open(svg_file_path, "w", encoding="utf-8") as f:
                    f.write(svg_content)

            # Speaker notes
            note_content = notes_by_slide.get(idx + 1) or ""
            if idx == 0 and material.active_learning_script and "---RATIONALE---" in material.active_learning_script:
                parts = material.active_learning_script.split("---RATIONALE---", 1)
                rationale_text = parts[1].strip()
                if rationale_text:
                    note_content = f"💡 PEDAGOGICAL RATIONALE:\n\n{rationale_text}\n\n{note_content}".strip()

            if note_content:
                note_file_path = os.path.join(notes_dir, f"{slide_name}.md")
                with open(note_file_path, "w", encoding="utf-8") as f:
                    f.write(note_content)

        # 4. Invoke PPT-Master svg_to_pptx.py
        ppt_master_script = _PPT_MASTER_SCRIPT_PATH
        if not os.path.isfile(ppt_master_script):
            raise HTTPException(
                status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
                detail=f"PPT-Master script không tìm thấy: {ppt_master_script}. "
                f"Hãy cấu hình biến môi trường PPT_MASTER_SCRIPT_PATH.",
            )
        output_pptx = os.path.join(exports_dir, "output.pptx")
        cmd = f'"{sys.executable}" "{ppt_master_script}" "{project_dir}" -o "{output_pptx}" --no-compat'

        # Execute command
        res = subprocess.run(cmd, shell=True, capture_output=True, text=True)
        if res.returncode != 0:
            print(f"PPT-Master generate error: {res.stderr}")
            raise HTTPException(
                status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
                detail=f"Lỗi PPT-Master CLI: {res.stderr or res.stdout}",
            )

        return FileResponse(
            output_pptx,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            filename=f"Bai_Giang_Chuong_{chapter_id}.pptx",
        )


def render_markdown_to_html(md_text: str) -> str:
    # Escape HTML characters safely
    html = md_text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
    # Title markers
    html = re.sub(r"^### (.*?)$", r"<h3 class='lp-h3'>\1</h3>", html, flags=re.MULTILINE)
    html = re.sub(r"^## (.*?)$", r"<h2 class='lp-h2'>\1</h2>", html, flags=re.MULTILINE)
    html = re.sub(r"^# (.*?)$", r"<h1 class='lp-h1'>\1</h1>", html, flags=re.MULTILINE)
    # Bold markers
    html = re.sub(r"\*\*(.*?)\*\*", r"<strong>\1</strong>", html)
    # List parsing
    lines = html.split("\n")
    in_list = False
    new_lines = []
    for line in lines:
        match = re.match(r"^[-*+•]\s*(.*)$", line.strip())
        if match:
            if not in_list:
                new_lines.append("<ul class='lp-ul'>")
                in_list = True
            new_lines.append(f"<li class='lp-li'>{match.group(1)}</li>")
        else:
            if in_list:
                new_lines.append("</ul>")
                in_list = False
            new_lines.append(line)
    if in_list:
        new_lines.append("</ul>")
    html = "\n".join(new_lines)
    # Paragraph mapping
    html = "\n".join(
        f"<p class='lp-p'>{line}</p>"
        if line.strip()
        and not line.strip().startswith("<h")
        and not line.strip().startswith("<u")
        and not line.strip().startswith("</u")
        and not line.strip().startswith("<l")
        else line
        for line in html.split("\n")
    )
    return html


@router.get("/chapters/{chapter_id}/export-lesson-plan", response_class=HTMLResponse)
def export_lesson_plan(chapter_id: int, current_user: User = Depends(get_current_user), db: Session = Depends(get_db)):
    # 1. Xác thực chương học và quyền sở hữu môn học
    chapter = (
        db.query(Chapter)
        .join(Course)
        .filter(Chapter.id == chapter_id, Course.user_id == current_user.id, Chapter.is_active)
        .first()
    )
    if not chapter:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND, detail="Chương học không tồn tại hoặc bạn không có quyền truy cập."
        )

    # 2. Lấy học liệu chương học
    material = db.query(ChapterMaterial).filter(ChapterMaterial.chapter_id == chapter_id).first()
    if not material or not material.active_learning_script:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="Chương học này chưa có nội dung kịch bản giảng dạy Active Learning.",
        )

    active_learning_script = material.active_learning_script or ""
    marker = "---RATIONALE---"
    rationale_html = ""

    if marker in active_learning_script:
        parts = active_learning_script.split(marker, 1)
        main_script = parts[0].strip()
        rationale_text = parts[1].strip()
        script_html = render_markdown_to_html(main_script)
        if rationale_text:
            rationale_html = f"""
            <div class="rationale-panel" style="margin-top: 30px; padding: 20px; background-color: #f0fff4; border-left: 4px solid #38a169; border-radius: 8px; border: 1px solid #c6f6d5;">
                <h3 style="margin-top: 0; color: #276749; font-size: 14px; font-weight: bold; text-transform: uppercase; letter-spacing: 0.5px; display: flex; align-items: center; gap: 6px; margin-bottom: 10px;">
                    💡 GIẢI TRÌNH SƯ PHẠM (PEDAGOGICAL RATIONALE)
                </h3>
                <div style="font-size: 13.5px; color: #2f855a; font-style: italic; line-height: 1.5;">
                    {render_markdown_to_html(rationale_text)}
                </div>
            </div>
            """
    else:
        script_html = render_markdown_to_html(active_learning_script)

    course_name = chapter.course.course_name
    chapter_title = chapter.title
    author_name = current_user.full_name or current_user.email or "Giảng viên"

    html_content = f"""<!DOCTYPE html>
<html>
<head>
    <meta charset="utf-8">
    <title>Giáo án tương tác - Chương: {chapter_title}</title>
    <style>
        body {{
            font-family: 'Inter', -apple-system, BlinkMacSystemFont, "Segoe UI", Roboto, sans-serif;
            color: #1a202c;
            line-height: 1.6;
            padding: 40px;
            max-width: 800px;
            margin: 0 auto;
            background-color: #ffffff;
        }}
        @media print {{
            body {{
                padding: 0;
                max-width: 100%;
                font-size: 12pt;
            }}
            .no-print {{
                display: none;
            }}
        }}
        .header-panel {{
            border-bottom: 2px solid #1a365d;
            padding-bottom: 20px;
            margin-bottom: 30px;
        }}
        .meta-grid {{
            display: grid;
            grid-template-columns: 1fr 1fr;
            gap: 15px;
            margin-top: 15px;
            font-size: 14px;
            background-color: #f7fafc;
            padding: 15px;
            border-radius: 8px;
            border: 1px solid #edf2f7;
        }}
        .meta-item {{
            margin-bottom: 5px;
        }}
        .meta-label {{
            font-weight: bold;
            color: #4a5568;
        }}
        .lp-h1 {{
            color: #1a365d;
            font-size: 24px;
            margin-top: 0;
            margin-bottom: 5px;
        }}
        .lp-h2 {{
            color: #2c5282;
            font-size: 18px;
            margin-top: 25px;
            margin-bottom: 12px;
            border-bottom: 1px solid #e2e8f0;
            padding-bottom: 5px;
        }}
        .lp-h3 {{
            color: #4a5568;
            font-size: 15px;
            margin-top: 18px;
            margin-bottom: 8px;
        }}
        .lp-p {{
            margin-bottom: 12px;
            text-align: justify;
        }}
        .lp-ul {{
            margin: 10px 0 15px 20px;
            padding: 0;
        }}
        .lp-li {{
            margin-bottom: 6px;
        }}
        .print-btn {{
            background-color: #1a365d;
            color: #ffffff;
            border: none;
            padding: 10px 20px;
            font-size: 14px;
            font-weight: bold;
            border-radius: 6px;
            cursor: pointer;
            transition: background 0.2s;
            margin-bottom: 20px;
        }}
        .print-btn:hover {{
            background-color: #2b6cb0;
        }}
    </style>
</head>
<body>
    <div class="no-print" style="text-align: right;">
        <button class="print-btn" onclick="window.print()">🖨️ In hoặc Lưu file PDF Giáo án</button>
    </div>
    <div class="header-panel">
        <h1 class="lp-h1">GIÁO ÁN HOẠT ĐỘNG SƯ PHẠM (LESSON PLAN)</h1>
        <div style="font-size: 14px; color: #718096; font-style: italic;">AI Lecture Assistant - Đảm bảo Chất lượng Đào tạo</div>
        <div class="meta-grid">
            <div class="meta-item"><span class="meta-label">Môn học:</span> {course_name}</div>
            <div class="meta-item"><span class="meta-label">Chương học:</span> {chapter_title}</div>
            <div class="meta-item"><span class="meta-label">Giảng viên:</span> {author_name}</div>
            <div class="meta-item"><span class="meta-label">Thiết kế Sư phạm:</span> Tương tác chủ động (Active Learning)</div>
        </div>
    </div>
    <div class="content-body">
        {script_html}
        {rationale_html}
    </div>
</body>
</html>"""
    return HTMLResponse(content=html_content, status_code=200)


class SlideItemPayload(BaseModel):
    type: str
    rawText: str | None = None  # noqa: N815
    bullet: bool | None = True


class SlidePayload(BaseModel):
    title: str
    layout: str
    items: list[SlideItemPayload]
    notes: str | None = None
    screenshot: str | None = None  # Full slide screenshot
    has_visual: bool | None = False
    visual_screenshot: str | None = None  # Base64 of diagram/chart/table


class ExportCanvasPayload(BaseModel):
    slides: list[SlidePayload]
    theme: str | None = "warm_academic"


@router.post("/chapters/{chapter_id}/export-pptx-canvas")
def export_chapter_pptx_canvas(
    chapter_id: int,
    payload: ExportCanvasPayload,
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db),
):
    # 1. Xác thực chương học và quyền sở hữu môn học
    chapter = (
        db.query(Chapter)
        .join(Course)
        .filter(Chapter.id == chapter_id, Course.user_id == current_user.id, Chapter.is_active)
        .first()
    )
    if not chapter:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND, detail="Chương học không tồn tại hoặc bạn không có quyền truy cập."
        )

    # 2. Setup presentation
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 3. Choose theme colors
    themes_colors = {
        "deep_space": {"bg": "#0A0A1A", "text": "#FFFFFF", "accent": "#00D2FF", "sub": "#8899BB"},
        "warm_academic": {"bg": "#FAF6EE", "text": "#1A202C", "accent": "#8C6239", "sub": "#5A6A80"},
        "mint_techno": {"bg": "#0B132B", "text": "#FFFFFF", "accent": "#1DE9B6", "sub": "#B2DFDB"},
        "sunset_crimson": {"bg": "#1A0813", "text": "#FFFFFF", "accent": "#FF5252", "sub": "#FF8A80"},
        "mckinsey_consulting": {"bg": "#041E42", "text": "#FFFFFF", "accent": "#00A3A6", "sub": "#80CBC4"},
    }

    theme_name = payload.theme or "warm_academic"
    colors = themes_colors.get(theme_name, themes_colors["warm_academic"])

    def hex_to_rgb(hex_str):
        hex_str = hex_str.lstrip("#")
        return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))

    blank_layout = prs.slide_layouts[6]  # Blank layout

    for s_idx, s in enumerate(payload.slides):
        slide = prs.slides.add_slide(blank_layout)

        # Set slide background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = hex_to_rgb(colors["bg"])

        # Add slide title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(1.0))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = s.title
        p_title.font.name = "Arial"
        p_title.font.size = Pt(32)
        p_title.font.bold = True
        p_title.font.color.rgb = hex_to_rgb(colors["accent"])

        # Determine if slide has a visual element (Mermaid, Chart, Table)
        has_visual = s.has_visual and s.visual_screenshot

        # Write slide items (bullet points)
        if has_visual:
            # Dựng 2 cột
            # Cột trái: Văn bản (Bullet points)
            text_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.8))
            tf = text_box.text_frame
            tf.word_wrap = True

            # Cột phải: Chèn ảnh chụp cấu phần trực quan
            try:
                img_data = s.visual_screenshot
                if "," in img_data:
                    img_data = img_data.split(",")[1]
                image_bytes = base64.b64decode(img_data)
                image_stream = BytesIO(image_bytes)

                # Insert picture with right position (centered in the right column)
                slide.shapes.add_picture(image_stream, Inches(6.8), Inches(1.8), width=Inches(5.7))
            except Exception as e:
                print(f"Error inserting visual screenshot on slide {s_idx + 1}: {e}")
                # Fallback: make text box wider
                text_box.width = Inches(11.7)
        else:
            # Slide chỉ có văn bản (Text-only)
            text_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.8))
            tf = text_box.text_frame
            tf.word_wrap = True

        # Add items to text box
        text_items = [it for it in s.items if it.type == "text" or it.rawText]
        for item_idx, item in enumerate(text_items):
            raw_text = item.rawText or ""
            clean_text = raw_text.replace("**", "").replace("* ", "").strip()
            if not clean_text:
                continue

            if item_idx == 0 and len(tf.paragraphs) > 0 and not tf.paragraphs[0].text:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            p.text = "• " + clean_text if item.bullet else clean_text
            p.font.name = "Arial"
            p.font.size = Pt(18 if len(text_items) > 5 else 20)
            p.font.color.rgb = hex_to_rgb(colors["text"])
            # Space after paragraphs
            p.space_after = Pt(12)

        # Add slide notes
        if s.notes:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = s.notes

    # Save presentation to a temporary path
    project_dir = os.path.join("temp", f"ppt_canvas_{chapter_id}")
    os.makedirs(project_dir, exist_ok=True)
    output_pptx = os.path.join(project_dir, f"chapter_{chapter_id}_canvas.pptx")
    prs.save(output_pptx)

    return FileResponse(
        output_pptx,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename=f"Bai_Giang_Chuong_{chapter_id}.pptx",
    )


# --- API EXPORT ZIP COURSE PACKAGE ---


def format_to_gift(question_text: str, options: list, correct_answer: str) -> str:
    """Formats a question into Moodle/Canvas GIFT format."""
    def escape_gift(text):
        for char in ['{', '}', '~', '=', '#', ':', '/']:
            text = text.replace(char, '\\' + char)
        return text

    q_text = escape_gift(question_text)
    gift_options = []

    # Normal case: correct answer matches one of options
    for opt in options:
        opt_escaped = escape_gift(opt)
        if opt.strip().lower() == correct_answer.strip().lower():
            gift_options.append(f"={opt_escaped}")
        else:
            gift_options.append(f"~{opt_escaped}")

    # Fallback if correct_answer was letter matching index (e.g., 'A', 'B', 'C', 'D')
    has_correct = any(o.startswith('=') for o in gift_options)
    if not has_correct and correct_answer in ["A", "B", "C", "D", "a", "b", "c", "d"]:
        idx = ord(correct_answer.upper()) - ord('A')
        if idx < len(gift_options):
            gift_options[idx] = "=" + gift_options[idx].lstrip('~')
            has_correct = True

    # Ultimate fallback: first item
    if not has_correct and gift_options:
        gift_options[0] = "=" + gift_options[0].lstrip('~')

    options_str = " ".join(gift_options)
    return f"{q_text} {{{options_str}}}"


def sanitize_filename(name: str) -> str:
    """Converts name to standard ASCII alphanumeric string for folder/file names."""
    # Normalize unicode to separate diacritics
    normalized = unicodedata.normalize('NFKD', name)
    ascii_encoded = normalized.encode('ascii', 'ignore').decode('ascii')
    # Filter non-alphanumeric/spaces/hyphens
    filtered = re.sub(r'[^a-zA-Z0-9_\-\s]', '', ascii_encoded)
    # Strip and replace spaces/underscores with a single underscore
    sanitized = re.sub(r'[\s_]+', '_', filtered.strip())
    return sanitized or "Chapter"


@router.get("/{course_id}/export-zip")
def export_course_zip(
    course_id: int,
    background_tasks: BackgroundTasks,
    organization_style: str = "by_chapter",
    theme: str = "warm_academic",
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    # 1. Xác thực môn học
    course = db.query(Course).filter(Course.id == course_id, Course.user_id == current_user.id).first()
    if not course:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND, detail="Môn học không tồn tại hoặc bạn không có quyền truy cập."
        )

    # 2. Lấy dữ liệu chương học, câu hỏi, CLO
    chapters = (
        db.query(Chapter)
        .filter(Chapter.course_id == course_id, Chapter.is_active)
        .order_by(Chapter.sort_order)
        .all()
    )
    questions = db.query(Question).filter(Question.course_id == course_id, Question.is_active).all()
    clos = db.query(CLO).filter(CLO.course_id == course_id).all()

    # 3. Tạo thư mục tạm để đóng gói
    temp_dir = tempfile.mkdtemp(prefix=f"course_export_{course_id}_")

    try:
        # A. Đề cương Syllabus.md
        syllabus_path = os.path.join(temp_dir, "Syllabus.md")
        with open(syllabus_path, "w", encoding="utf-8") as f:
            f.write(f"# ĐỀ CƯƠNG CHI TIẾT MÔN HỌC: {course.course_name.upper()}\n")
            f.write(f"Mã môn học: {course.course_code}\n")
            f.write(f"Giảng viên biên soạn: {current_user.full_name or current_user.email}\n\n")
            f.write("## 🎯 Chuẩn đầu ra môn học (CLOs)\n")
            if not clos:
                f.write("* Chưa có chuẩn đầu ra nào được cấu hình.\n")
            else:
                for c in clos:
                    f.write(f"- **{c.clo_code}** (Mức Bloom {c.bloom_level}): {c.description}\n")
            f.write("\n## 📚 Giáo trình & Tài liệu tham khảo\n")
            f.write(f"- Giáo trình bắt buộc: {course.required_textbooks or 'N/A'}\n")
            f.write(f"- Tài liệu tham khảo: {course.recommended_readings or 'N/A'}\n")

        # B. Báo cáo Ma trận Coverage Bloom x CLO
        matrix_path = os.path.join(temp_dir, "Matrix_Coverage.md")
        with open(matrix_path, "w", encoding="utf-8") as f:
            f.write("# MA TRẬN ĐỘ PHỦ CHẤT LƯỢNG (CLO x BLOOM LEVEL)\n")
            f.write(f"Môn học: {course.course_name} ({course.course_code})\n\n")
            f.write("| Chuẩn đầu ra (CLO) | Mức Bloom | Số lượng Câu hỏi | Slide bài giảng | Mô tả CLO |\n")
            f.write("| :--- | :---: | :---: | :---: | :--- |\n")

            for c in clos:
                q_count = len([q for q in questions if q.clo_id == c.id])
                # Check how many chapter materials mention this CLO
                slide_count = 0
                for ch in chapters:
                    mat = db.query(ChapterMaterial).filter(ChapterMaterial.chapter_id == ch.id).first()
                    if mat and mat.slide_content and c.clo_code in mat.slide_content:
                        slide_count += 1
                f.write(f"| {c.clo_code} | Mức {c.bloom_level} | {q_count} câu | {slide_count} chương | {c.description} |\n")

        # C. Đóng gói các chương theo cấu trúc thư mục mong muốn
        if organization_style == "by_type":
            # Tạo thư mục phân loại theo loại học liệu
            storyboard_dir = os.path.join(temp_dir, "Storyboards")
            slides_dir = os.path.join(temp_dir, "Slides")
            quizzes_dir = os.path.join(temp_dir, "Quizzes")
            os.makedirs(storyboard_dir, exist_ok=True)
            os.makedirs(slides_dir, exist_ok=True)
            os.makedirs(quizzes_dir, exist_ok=True)

        for idx, ch in enumerate(chapters):
            ch_num = idx + 1
            sanitized_title = sanitize_filename(ch.title)
            ch_folder_name = f"Chapter_{ch_num:02d}_{sanitized_title}"

            # Khởi tạo đường dẫn lưu file của chương này
            if organization_style == "by_type":
                storyboard_dest_dir = storyboard_dir
                slides_dest_dir = slides_dir
                quizzes_dest_dir = quizzes_dir
                file_prefix = f"Chapter_{ch_num:02d}_"
            else: # default: by_chapter
                ch_path = os.path.join(temp_dir, "Chapters", ch_folder_name)
                os.makedirs(ch_path, exist_ok=True)
                storyboard_dest_dir = ch_path
                slides_dest_dir = ch_path
                quizzes_dest_dir = ch_path
                file_prefix = ""

            material = db.query(ChapterMaterial).filter(ChapterMaterial.chapter_id == ch.id).first()

            # i. Kịch bản Active Learning (Storyboard)
            if material and material.active_learning_script:
                sb_path = os.path.join(storyboard_dest_dir, f"{file_prefix}Storyboard.md")
                with open(sb_path, "w", encoding="utf-8") as f:
                    f.write(f"# GIÁO ÁN ACTIVE LEARNING: {ch.title.upper()}\n\n")
                    f.write(material.active_learning_script)

            # ii. Slide PPTX (Fault-tolerant)
            if material and material.slide_content:
                # Lưu file slide nguồn Markdown trước làm dự phòng
                src_slide_path = os.path.join(slides_dest_dir, f"{file_prefix}Slides_Source.md")
                with open(src_slide_path, "w", encoding="utf-8") as f:
                    f.write(material.slide_content)

                # Cố gắng xuất slide PPTX
                try:
                    # Gọi trực tiếp helper logic của export_chapter_pptx nhưng không trả Response
                    output_pptx = export_chapter_pptx(chapter_id=ch.id, theme=theme, current_user=current_user, db=db)
                    # export_chapter_pptx trả về FileResponse, ta có thể lấy path từ nó
                    if hasattr(output_pptx, "path") and os.path.exists(output_pptx.path):
                        dest_pptx = os.path.join(slides_dest_dir, f"{file_prefix}Slide_Presentation.pptx")
                        shutil.copy(output_pptx.path, dest_pptx)
                except Exception as e:
                    # Ghi nhận lỗi và tiếp tục đóng gói
                    err_path = os.path.join(slides_dest_dir, f"{file_prefix}ERROR_Slides_Generation.txt")
                    with open(err_path, "w", encoding="utf-8") as f:
                        f.write(f"Lỗi sinh PowerPoint cho chương này:\n{str(e)}")

            # iii. Câu hỏi thi (Quiz - Markdown & GIFT format)
            ch_questions = [q for q in questions if q.chapter_id == ch.id]
            if ch_questions:
                # Markdown format quiz
                quiz_md_path = os.path.join(quizzes_dest_dir, f"{file_prefix}Quiz_Questions.md")
                with open(quiz_md_path, "w", encoding="utf-8") as f:
                    f.write(f"# NGÂN HÀNG CÂU HỎI CHƯƠNG {ch_num}: {ch.title}\n\n")
                    for q_idx, q in enumerate(ch_questions):
                        f.write(f"Câu {q_idx + 1}: {q.question_text}\n")
                        opts = []
                        try:
                            opts = json.loads(q.options_json) if q.options_json else []
                        except Exception:
                            pass
                        labels = ["A", "B", "C", "D"]
                        for opt_i, opt in enumerate(opts):
                            if opt_i < len(labels):
                                f.write(f"  {labels[opt_i]}. {opt}\n")
                        f.write(f"\n  * Đáp án đúng: {q.correct_answer}\n")
                        f.write(f"  * Cấp độ Bloom: Mức {q.bloom_level}\n\n")

                # GIFT format quiz (for Canvas/Moodle import)
                quiz_gift_path = os.path.join(quizzes_dest_dir, f"{file_prefix}Quiz_Questions.gift")
                with open(quiz_gift_path, "w", encoding="utf-8") as f:
                    f.write(f"// Ngân hàng câu hỏi trắc nghiệm Chương {ch_num}\n\n")
                    for q in ch_questions:
                        opts = []
                        try:
                            opts = json.loads(q.options_json) if q.options_json else []
                        except Exception:
                            pass
                        gift_str = format_to_gift(q.question_text, opts, q.correct_answer)
                        f.write(gift_str + "\n\n")

        # 4. Nén thư mục tạm thành file ZIP
        zip_temp_dir = tempfile.mkdtemp(prefix="course_zip_")
        zip_file_path = os.path.join(zip_temp_dir, f"Course_Package_{course.course_code}.zip")

        with zipfile.ZipFile(zip_file_path, "w", zipfile.ZIP_DEFLATED) as zip_file:
            for root, dirs, files in os.walk(temp_dir):
                for file in files:
                    file_path = os.path.join(root, file)
                    rel_path = os.path.relpath(file_path, temp_dir)
                    zip_file.write(file_path, rel_path)

        # 5. Đăng ký background task dọn dẹp các thư mục tạm sau khi truyền file xong
        def cleanup_temp_directories():
            try:
                shutil.rmtree(temp_dir)
                shutil.rmtree(zip_temp_dir)
            except Exception as e:
                print(f"Lỗi khi dọn dẹp thư mục tạm ZIP export: {e}")

        background_tasks.add_task(cleanup_temp_directories)

        # 6. Trả về file ZIP
        return FileResponse(
            zip_file_path,
            media_type="application/zip",
            filename=f"Course_Package_{course.course_code}.zip"
        )

    except Exception as outer_err:
        # Nếu lỗi trong quá trình xử lý, cố gắng dọn dẹp thư mục tạm
        try:
            shutil.rmtree(temp_dir)
        except Exception:
            pass
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"Lỗi đóng gói file ZIP: {str(outer_err)}"
        )


