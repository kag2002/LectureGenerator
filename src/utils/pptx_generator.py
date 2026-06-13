from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# Theme colors
themes_bg = {
    "deep_space": "0A0A1A",
    "warm_academic": "FAF6EE",
    "mint_techno": "0B132B",
    "sunset_crimson": "1A0813",
    "mckinsey_consulting": "041E42",
}
themes_accent = {
    "deep_space": "00D2FF",
    "warm_academic": "8C6239",
    "mint_techno": "1DE9B6",
    "sunset_crimson": "FF5252",
    "mckinsey_consulting": "00A3A6",
}
themes_text = {
    "deep_space": "E2E8F0",
    "warm_academic": "2D3748",
    "mint_techno": "E2E8F0",
    "sunset_crimson": "F8FAFC",
    "mckinsey_consulting": "FFFFFF",
}
themes_sub = {
    "deep_space": "8899BB",
    "warm_academic": "5A6A80",
    "mint_techno": "B2DFDB",
    "sunset_crimson": "FF8A80",
    "mckinsey_consulting": "80CBC4",
}


def generate_pptx_from_slidej_json(slidej_json: dict, output_path: str, theme_name: str = "warm_academic"):
    """
    Fallback: Sinh file PPTX trực tiếp từ cấu trúc SlideJ JSON bằng python-pptx,
    không phụ thuộc SlideJ CLI.
    """
    width_inch = slidej_json.get("width", 13.333)
    height_inch = slidej_json.get("height", 7.5)

    prs = Presentation()
    prs.slide_width = Inches(width_inch)
    prs.slide_height = Inches(height_inch)

    # Blank layout
    blank_layout = prs.slide_layouts[6]

    def hex_to_rgb(hex_str: str) -> RGBColor:
        hex_str = hex_str.strip().lstrip("#")
        if len(hex_str) == 6:
            return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))
        return RGBColor(255, 255, 255)

    def set_fill_solid(shape, color_hex: str):
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(color_hex)

    def add_text_box(slide, x, y, w, h, text, font_size=18, color="FFFFFF", bold=False, align="left"):
        txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_top = Pt(4)
        tf.margin_bottom = Pt(4)
        tf.margin_left = Pt(8)
        tf.margin_right = Pt(8)

        # Handle text with **bold** markers
        p = tf.paragraphs[0]
        if align == "center":
            p.alignment = PP_ALIGN.CENTER
        elif align == "right":
            p.alignment = PP_ALIGN.RIGHT
        else:
            p.alignment = PP_ALIGN.LEFT

        parts = str(text).split("**")
        for idx, part in enumerate(parts):
            if not part:
                continue
            is_bold_part = idx % 2 == 1
            run = p.add_run()
            run.text = part
            run.font.size = Pt(font_size)
            run.font.bold = bold or is_bold_part
            run.font.color.rgb = hex_to_rgb(color)

        return txBox

    def add_bg_rect(slide, color_hex: str):
        bg_shape = slide.shapes.add_shape(
            1,  # MSO_SHAPE.RECTANGLE
            Inches(0),
            Inches(0),
            Inches(width_inch),
            Inches(height_inch),
        )
        set_fill_solid(bg_shape, color_hex)
        bg_shape.line.fill.background()

    bg_color = themes_bg.get(theme_name, "0A0A1A")
    accent_color = themes_accent.get(theme_name, "00D2FF")
    text_color = themes_text.get(theme_name, "E2E8F0")
    themes_sub.get(theme_name, "8899BB")

    slides_data = slidej_json.get("slides", [])

    for s_idx, slide_data in enumerate(slides_data):
        slide = prs.slides.add_slide(blank_layout)

        # Background
        bg = slide_data.get("background", bg_color)
        if isinstance(bg, str):
            add_bg_rect(slide, bg)
        elif isinstance(bg, dict):
            # Use first stop color for solid fill fallback
            stops = bg.get("stops", [])
            bg_hex = stops[0].get("color", bg_color) if stops else bg_color
            add_bg_rect(slide, bg_hex)

        # Elements
        for elem in slide_data.get("elements", []):
            elem_type = elem.get("type", "")
            pos = elem.get("position", {})
            x = pos.get("x", 0)
            y = pos.get("y", 0)
            w = pos.get("w", 4)
            h = pos.get("h", 1)

            if elem_type == "text":
                text_content = elem.get("text", "")
                font_size = elem.get("fontSize", 18)
                color = elem.get("color", text_color)
                bold = elem.get("bold", False)
                align = elem.get("align", "left")

                # Handle paragraphs list (SlideJ format with runs)
                if isinstance(text_content, list):
                    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
                    tf = txBox.text_frame
                    tf.word_wrap = True
                    tf.margin_top = Pt(4)
                    tf.margin_bottom = Pt(4)
                    tf.margin_left = Pt(8)
                    tf.margin_right = Pt(8)

                    for p_idx, para_data in enumerate(text_content):
                        if p_idx == 0:
                            p = tf.paragraphs[0]
                        else:
                            p = tf.add_paragraph()

                        if align == "center":
                            p.alignment = PP_ALIGN.CENTER
                        elif align == "right":
                            p.alignment = PP_ALIGN.RIGHT
                        else:
                            p.alignment = PP_ALIGN.LEFT

                        is_bullet = para_data.get("bullet", False)
                        if is_bullet:
                            p.level = 0

                        runs = para_data.get("runs", [])
                        for run_data in runs:
                            run = p.add_run()
                            run.text = run_data.get("text", "")
                            run.font.size = Pt(run_data.get("fontSize", font_size))
                            run.font.bold = run_data.get("bold", False)
                            run_color = run_data.get("color", color)
                            run.font.color.rgb = hex_to_rgb(run_color)
                else:
                    add_text_box(slide, x, y, w, h, text_content, font_size, color, bold, align)

            elif elem_type == "shape":
                shape_type = elem.get("shapeType", "rect")
                from pptx.enum.shapes import MSO_SHAPE

                shape_enum = MSO_SHAPE.RECTANGLE
                if shape_type == "roundRect":
                    shape_enum = MSO_SHAPE.ROUNDED_RECTANGLE
                elif shape_type == "ellipse":
                    shape_enum = MSO_SHAPE.OVAL

                shape = slide.shapes.add_shape(shape_enum, Inches(x), Inches(y), Inches(w), Inches(h))

                fill = elem.get("fill", None)
                if fill:
                    if isinstance(fill, str):
                        set_fill_solid(shape, fill)
                    elif isinstance(fill, dict):
                        stops = fill.get("stops", [])
                        if stops:
                            set_fill_solid(shape, stops[0].get("color", "333333"))

                line = elem.get("line", None)
                if line:
                    line_color = line.get("color", accent_color)
                    line_width = line.get("width", 1)
                    shape.line.color.rgb = hex_to_rgb(line_color)
                    shape.line.width = Pt(line_width)
                else:
                    shape.line.fill.background()

                # Text inside shape
                shape_text = elem.get("text", None)
                if shape_text and isinstance(shape_text, list):
                    tf = shape.text_frame
                    tf.word_wrap = True
                    margin_inch = elem.get("margin", 0.15)
                    tf.margin_top = Inches(margin_inch)
                    tf.margin_bottom = Inches(margin_inch)
                    tf.margin_left = Inches(margin_inch)
                    tf.margin_right = Inches(margin_inch)

                    vert_align = elem.get("vertAlign", "top")
                    if vert_align == "middle":
                        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

                    for p_idx, para_data in enumerate(shape_text):
                        if p_idx == 0:
                            p = tf.paragraphs[0]
                        else:
                            p = tf.add_paragraph()

                        para_align = para_data.get("align", elem.get("align", "left"))
                        if para_align == "center":
                            p.alignment = PP_ALIGN.CENTER
                        else:
                            p.alignment = PP_ALIGN.LEFT

                        for run_data in para_data.get("runs", []):
                            run = p.add_run()
                            run.text = run_data.get("text", "")
                            run.font.size = Pt(run_data.get("fontSize", 14))
                            run.font.bold = run_data.get("bold", False)
                            run_color = run_data.get("color", text_color)
                            run.font.color.rgb = hex_to_rgb(run_color)

            elif elem_type == "table":
                rows = elem.get("rows", [])
                if not rows:
                    continue
                num_rows = len(rows)
                num_cols = max(len(r) if isinstance(r, list) else 0 for r in rows) if rows else 1

                table_shape = slide.shapes.add_table(num_rows, num_cols, Inches(x), Inches(y), Inches(w), Inches(h))
                table = table_shape.table

                for r_idx, row in enumerate(rows):
                    if not isinstance(row, list):
                        continue
                    for c_idx, cell_data in enumerate(row):
                        if c_idx >= num_cols:
                            break
                        cell = table.cell(r_idx, c_idx)
                        if isinstance(cell_data, dict):
                            cell.text = cell_data.get("text", "")
                            for p in cell.text_frame.paragraphs:
                                for run in p.runs:
                                    run.font.size = Pt(elem.get("fontSize", 12))
                                    if cell_data.get("bold"):
                                        run.font.bold = True
                                    if cell_data.get("color"):
                                        run.font.color.rgb = hex_to_rgb(cell_data["color"])
                            if cell_data.get("fill"):
                                cell.fill.solid()
                                cell.fill.fore_color.rgb = hex_to_rgb(cell_data["fill"])
                        else:
                            cell.text = str(cell_data) if cell_data else ""
                            for p in cell.text_frame.paragraphs:
                                for run in p.runs:
                                    run.font.size = Pt(elem.get("fontSize", 12))
                                    run.font.color.rgb = hex_to_rgb(elem.get("color", text_color))

    prs.save(output_path)
    print(f"[INFO] python-pptx fallback: Generated {len(slides_data)} slides → {output_path}")
