import re
import xml.etree.ElementTree as ET

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.util import Inches, Pt


def draw_svg_on_slide(slide, svg_content: str):
    """
    Dịch và vẽ nội dung SVG XML thành các vector shape gốc trên slide PowerPoint.
    """
    try:
        clean_svg = re.sub(r'\sxmlns(:\w+)?="[^"]+"', "", svg_content)
        root = ET.fromstring(clean_svg)
    except Exception as e:
        print(f"[ERROR] Failed to parse SVG XML: {e}")
        return

    viewbox = root.get("viewBox")
    if viewbox:
        parts = [float(p) for p in viewbox.split()]
        if len(parts) == 4:
            svg_w = parts[2] - parts[0]
            svg_h = parts[3] - parts[1]
        else:
            svg_w = float(root.get("width", "800").replace("px", ""))
            svg_h = float(root.get("height", "450").replace("px", ""))
    else:
        svg_w = float(root.get("width", "800").replace("px", ""))
        svg_h = float(root.get("height", "450").replace("px", ""))

    if svg_w <= 0 or svg_h <= 0:
        return

    box_x = Inches(1.0)
    box_y = Inches(1.8)
    box_w = Inches(11.33)
    box_h = Inches(4.5)

    scale_x = box_w / svg_w
    scale_y = box_h / svg_h
    scale = min(scale_x, scale_y)

    offset_x = box_x + (box_w - svg_w * scale) / 2
    offset_y = box_y + (box_h - svg_h * scale) / 2

    def transform_coords(x, y):
        return int(offset_x + x * scale), int(offset_y + y * scale)

    def transform_dim(d):
        return int(d * scale)

    def hex_to_rgb(hex_str):
        if not hex_str or hex_str == "none":
            return None
        hex_str = hex_str.strip().lstrip("#")
        if len(hex_str) == 3:
            hex_str = "".join(c * 2 for c in hex_str)
        if len(hex_str) == 6:
            try:
                return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))
            except ValueError:
                return None
        colors = {
            "white": RGBColor(255, 255, 255),
            "black": RGBColor(0, 0, 0),
            "red": RGBColor(255, 0, 0),
            "green": RGBColor(0, 255, 0),
            "blue": RGBColor(0, 0, 255),
            "yellow": RGBColor(255, 255, 0),
            "cyan": RGBColor(0, 255, 255),
            "magenta": RGBColor(255, 0, 255),
            "teal": RGBColor(0, 128, 128),
        }
        return colors.get(hex_str.lower(), None)

    for elem in root.iter():
        tag = elem.tag.split("}")[-1].lower()

        if tag == "rect":
            x = float(elem.get("x", "0"))
            y = float(elem.get("y", "0"))
            w = float(elem.get("width", "0"))
            h = float(elem.get("height", "0"))
            fill = elem.get("fill", "none")
            stroke = elem.get("stroke", "none")
            stroke_width = float(elem.get("stroke-width", "1").replace("px", ""))

            l, t = transform_coords(x, y)
            width = transform_dim(w)
            height = transform_dim(h)

            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, width, height)

            rgb_fill = hex_to_rgb(fill)
            if rgb_fill:
                shape.fill.solid()
                shape.fill.fore_color.rgb = rgb_fill
            else:
                shape.fill.background()

            rgb_stroke = hex_to_rgb(stroke)
            if rgb_stroke:
                shape.line.color.rgb = rgb_stroke
                shape.line.width = Pt(stroke_width)
            else:
                shape.line.fill.background()

        elif tag == "circle":
            cx = float(elem.get("cx", "0"))
            cy = float(elem.get("cy", "0"))
            r = float(elem.get("r", "0"))
            fill = elem.get("fill", "none")
            stroke = elem.get("stroke", "none")
            stroke_width = float(elem.get("stroke-width", "1").replace("px", ""))

            l, t = transform_coords(cx - r, cy - r)
            width = transform_dim(2 * r)
            height = transform_dim(2 * r)

            shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, width, height)

            rgb_fill = hex_to_rgb(fill)
            if rgb_fill:
                shape.fill.solid()
                shape.fill.fore_color.rgb = rgb_fill
            else:
                shape.fill.background()

            rgb_stroke = hex_to_rgb(stroke)
            if rgb_stroke:
                shape.line.color.rgb = rgb_stroke
                shape.line.width = Pt(stroke_width)
            else:
                shape.line.fill.background()

        elif tag == "line":
            x1 = float(elem.get("x1", "0"))
            y1 = float(elem.get("y1", "0"))
            x2 = float(elem.get("x2", "0"))
            y2 = float(elem.get("y2", "0"))
            stroke = elem.get("stroke", "none")
            stroke_width = float(elem.get("stroke-width", "1").replace("px", ""))

            sx, sy = transform_coords(x1, y1)
            ex, ey = transform_coords(x2, y2)

            connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, sx, sy, ex, ey)

            rgb_stroke = hex_to_rgb(stroke)
            if rgb_stroke:
                connector.line.color.rgb = rgb_stroke
                connector.line.width = Pt(stroke_width)

        elif tag == "text":
            x = float(elem.get("x", "0"))
            y = float(elem.get("y", "0"))
            font_size = float(elem.get("font-size", "14").replace("px", ""))
            fill = elem.get("fill", "#FFFFFF")
            text_content = elem.text or ""
            text_anchor = elem.get("text-anchor", "start")

            l, t = transform_coords(x, y - font_size)
            w = transform_dim(150)
            h = transform_dim(font_size * 1.5)

            txBox = slide.shapes.add_textbox(l, t, w, h)
            tf = txBox.text_frame
            tf.word_wrap = False
            tf.margin_top = 0
            tf.margin_bottom = 0
            tf.margin_left = 0
            tf.margin_right = 0

            p = tf.paragraphs[0]
            p.text = text_content

            rgb_fill = hex_to_rgb(fill)
            if rgb_fill:
                p.font.color.rgb = rgb_fill

            p.font.size = Pt(max(8, int(font_size * scale / Inches(1) * 72)))

            if text_anchor == "middle":
                from pptx.enum.text import PP_ALIGN

                p.alignment = PP_ALIGN.CENTER
